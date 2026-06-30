import streamlit as st
from google import genai
from google.genai import types
import json
import io
import time  # Required for managing retry delays
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# --- STEP 1: UI setup & Key configuration ---
st.set_page_config(page_title="OmniClutch AI", page_icon="🚨", layout="centered")

# Place this right after st.set_page_config()
def local_css():
    st.markdown(
        """
        <style>
        /* Main application background and font smoothing */
        .stApp {
            background-color: #0b192c;
            color: #ffffff;
        }
        
        /* Premium custom metric cards */
        .metric-card {
            background: rgba(255, 255, 255, 0.03);
            border: 1px solid rgba(255, 255, 255, 0.1);
            padding: 20px;
            border-radius: 12px;
            box-shadow: 0 4px 15px rgba(0,0,0,0.2);
            margin-bottom: 15px;
        }
        
        .metric-title {
            color: #ea5455;
            font-size: 14px;
            font-weight: bold;
            text-transform: uppercase;
            letter-spacing: 1px;
        }
        
        .metric-value {
            font-size: 24px;
            font-weight: 700;
            margin-top: 5px;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

# Run the style injector
local_css()


st.sidebar.title("Configuration")
GEMINI_KEY = st.sidebar.text_input("🔑 Enter Gemini API Key", type="password")

if not GEMINI_KEY:
    st.sidebar.warning("Please enter your Gemini API key.")
    st.info("💡 Paste your Gemini API key in the sidebar to activate the engine.")
    st.stop()

# Initialize the modern client
client = genai.Client(api_key=GEMINI_KEY)

# --- STEP 2: Slide Generation Logic ---
def generate_slide_deck(slide_data_list):
    prs = Presentation()
    NAVY = RGBColor(11, 25, 44)
    CORAL = RGBColor(234, 84, 85)
    WHITE = RGBColor(255, 255, 255)
    
    for slide_data in slide_data_list:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY
        
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.75), Inches(8.5), Inches(1))
        p_title = title_box.text_frame.paragraphs[0]
        p_title.text = slide_data.get("title", "Overview")
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = CORAL
        
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(8.5), Inches(4.5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        for idx, bullet_text in enumerate(slide_data.get("bullets", [])):
            p = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
            p.text = f"• {bullet_text}"
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.space_after = Pt(12)
            
    binary_stream = io.BytesIO()
    prs.save(binary_stream)
    binary_stream.seek(0)
    return binary_stream

# --- STEP 3: Main Dashboard Layout ---
# --- STEP 3: Premium UI & Main Dashboard Layout ---

# 1. Inject Custom CSS for the glowing download button
st.markdown("""
    <style>
    /* Premium Hover Effect for the Download Button */
    div[data-testid="stDownloadButton"] button {
        background-color: #ea5455;
        color: white;
        border-radius: 8px;
        border: none;
        font-size: 18px;
        font-weight: bold;
        transition: all 0.3s ease-in-out;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
    }
    div[data-testid="stDownloadButton"] button:hover {
        background-color: #ff6b6b;
        box-shadow: 0 6px 20px rgba(234, 84, 85, 0.5);
        transform: translateY(-2px);
        color: white;
    }
    </style>
    """, unsafe_allow_html=True)

# 2. Custom Gradient Header
st.markdown("""
    <div style="text-align: center; padding: 1rem 0; border-bottom: 1px solid rgba(255,255,255,0.1); margin-bottom: 2rem;">
        <h1 style="font-size: 3.5rem; margin-bottom: 0; background: -webkit-linear-gradient(45deg, #ea5455, #ffb26b); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">
            🚨 OmniClutch AI
        </h1>
        <p style="color: #a0aec0; font-size: 1.2rem; margin-top: 0.5rem; letter-spacing: 1px;">
            Turn Last-Minute Panic into Executive Deliverables.
        </p>
    </div>
    """, unsafe_allow_html=True)

# 3. Main Interface
audio_file = st.audio_input("Record your last-minute panic description:")

if st.button("🚀 Run OmniClutch Engine", use_container_width=True):
    if not audio_file:
        st.warning("Please record an audio file first.")
    else:
        try:
            with st.spinner("⚡ Gemini processing audio spectrum & structuring layouts..."):
                # [KEEP YOUR EXISTING GEMINI API CALL LOGIC HERE]
                audio_bytes = audio_file.read()
                audio_part = types.Part.from_bytes(data=audio_bytes, mime_type="audio/wav")
                
                prompt = """
                You are a crisis management presentation assistant. Analyze the attached audio panic note.
                Provide a structured response. You MUST reply ONLY with a valid JSON object matching this schema:
                {
                    "speech_script": "A beautifully drafted speech the user can read out loud right now.",
                    "slides": [
                        {"title": "Slide Title Here", "bullets": ["Critical point 1", "Critical point 2"]}
                    ]
                }
                """
                
                # Auto-Retry Loop
                max_retries = 3
                response = None
                for attempt in range(max_retries):
                    try:
                        response = client.models.generate_content(
                            model="gemini-2.5-flash",
                            contents=[prompt, audio_part],
                            config=types.GenerateContentConfig(response_mime_type="application/json")
                        )
                        break
                    except Exception as api_error:
                        if "503" in str(api_error) and attempt < max_retries - 1:
                            st.toast(f"⚠️ Server busy. Retrying (Attempt {attempt + 2}/{max_retries})...")
                            time.sleep(3)
                        else:
                            raise api_error
                
                ai_data = json.loads(response.text)
            
            # 4. Success Banner & Tabs
            st.success("🎉 Crisis Averted Successfully!")
            tab1, tab2 = st.tabs(["🗣️ Your Emergency Speech Script", "📊 Presentation Slide Matrix"])
            
            with tab1:
                st.write(ai_data.get("speech_script", "No script generated."))
                
            with tab2:
                slides = ai_data.get("slides", [])
                for s in slides:
                    st.markdown(f"**🎬 {s.get('title')}**")
                    for b in s.get("bullets", []):
                        st.write(f"- {b}")
                
                # The styled download button renders here
                st.markdown("<br>", unsafe_allow_html=True) # Adds a little breathing room
                pptx_output = generate_slide_deck(slides)
                st.download_button(
                    label="📥 Download Executive Slides (.pptx)",
                    data=pptx_output,
                    file_name="OmniClutch_Deck.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                
        except Exception as e:
            st.error(f"An error occurred during execution: {str(e)}")